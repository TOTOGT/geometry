#!/usr/bin/env python3
"""
Build the narration-script PDF from the deck's speaker notes.

The deck is the single source of truth: notes are read out of the .pptx, so the
script and the slides cannot drift apart. This PDF is what gets fed to an audio
generator — it contains the spoken text, which a slides-only export does not.
"""
from pptx import Presentation
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_LEFT
from reportlab.lib import colors
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                PageBreak, KeepTogether)
import re, html

SRC = 'HIST201-Wk01-Kepler.pptx'
OUT = 'HIST201-Wk01-Kepler-NARRATION.pdf'

NAVY = colors.HexColor('#1A2744')
GOLD = colors.HexColor('#8A6D १'.replace(' १', '2E'))  # placeholder-safe
GOLD = colors.HexColor('#8A6D2E')
INK  = colors.HexColor('#1A1A1A')
MUTE = colors.HexColor('#6B6B6B')

st_title = ParagraphStyle('t', fontName='Times-Bold', fontSize=26, leading=31,
                          textColor=NAVY, spaceAfter=6)
st_sub   = ParagraphStyle('s', fontName='Times-Italic', fontSize=13.5, leading=18,
                          textColor=GOLD, spaceAfter=20)
st_meta  = ParagraphStyle('m', fontName='Helvetica', fontSize=9.5, leading=15,
                          textColor=MUTE, spaceAfter=4)
st_part  = ParagraphStyle('p', fontName='Helvetica-Bold', fontSize=11, leading=15,
                          textColor=GOLD, spaceBefore=18, spaceAfter=14)
st_head  = ParagraphStyle('h', fontName='Helvetica-Bold', fontSize=12.5, leading=17,
                          textColor=NAVY, spaceBefore=16, spaceAfter=7)
st_body  = ParagraphStyle('b', fontName='Times-Roman', fontSize=12, leading=19.5,
                          textColor=INK, alignment=TA_LEFT, spaceAfter=9)
st_note  = ParagraphStyle('n', fontName='Helvetica-Oblique', fontSize=9.5, leading=14,
                          textColor=MUTE, spaceBefore=10, spaceAfter=10)

def esc(t):
    return html.escape(t).replace('—', '&#8212;').replace('’', '&#8217;')

# Section headings, in deck order. Named explicitly rather than scraped: the
# eyebrow line sits above the title on every slide and a "first short string"
# heuristic picks it up instead, which an audio generator would then read aloud.
TITLES = [
    'Kepler — the overreach, and what survived it',
    'A man who wanted the universe to be beautiful',
    'Model A — the five solids (Mysterium Cosmographicum, 1596)',
    'Model B — the third law (Harmonices Mundi, 1619)',
    'What you do this week — two columns, ten minutes',
    'Before you open the spreadsheet — post your prediction',
    'Part Two — what the numbers did',
    'Column one — the fitted one held',
    'Column two — the derived one did not',
    'The two columns, side by side',
    'What this establishes, and what it does not',
    'Carry this out of the week',
]

prs = Presentation(SRC)
slides = []
for i, sl in enumerate(prs.slides, start=1):
    title = TITLES[i-1] if i-1 < len(TITLES) else 'Slide %d' % i
    notes = ''
    if sl.has_notes_slide:
        notes = sl.notes_slide.notes_text_frame.text.strip()
    slides.append((i, title, notes))

assert len(prs.slides._sldIdLst) == len(TITLES), \
    'TITLES is out of sync with the deck (%d slides, %d titles)' % (
        len(prs.slides._sldIdLst), len(TITLES))

doc = SimpleDocTemplate(OUT, pagesize=LETTER,
                        leftMargin=1.05*inch, rightMargin=1.05*inch,
                        topMargin=0.95*inch, bottomMargin=0.9*inch,
                        title='HIST 201 Week 1 - Kepler - Lecture Narration',
                        author='Pablo Nogueira Grossi')

F = []
F.append(Paragraph('HIST&nbsp;201 &nbsp;&middot;&nbsp; Week&nbsp;1', st_meta))
F.append(Paragraph('Kepler: the overreach, and what survived it', st_title))
F.append(Paragraph('Lecture narration script &#8212; the spoken text for both clips', st_sub))
F.append(Paragraph(
    'This document contains the narration only. It is the source for the module audio: '
    'the words below are exactly what a listener hears, in order. '
    'The slides are a separate file and carry the visuals.', st_body))
F.append(Paragraph(
    'The lecture is deliberately split. Part One sets up the two competing models and '
    'stops before any result, because students must post a prediction before they compute. '
    'Part Two is released only after the computation deadline. Nothing in Part One reveals '
    'an outcome &#8212; that separation is a requirement, not a stylistic choice, and any '
    'edit to this script must preserve it.', st_body))
F.append(Paragraph(
    'Claim tags are spoken aloud where they appear (VERIFIED, MODEL, OPEN). They are load-bearing. '
    'A generated summary that flattens them into settled fact is wrong, and in this course '
    'that is itself an assessable error.', st_note))

part = None
for num, title, notes in slides:
    if not notes:
        continue
    new_part = 'PART TWO' if num >= 7 else 'PART ONE'
    if new_part != part:
        part = new_part
        F.append(PageBreak())
        label = ('PART ONE &#8212; released Monday. Watch before posting your prediction.'
                 if part == 'PART ONE' else
                 'PART TWO &#8212; released after the computation deadline.')
        F.append(Paragraph(label, st_part))
    F.append(Paragraph('%d. %s' % (num, esc(title)), st_head))
    for para in [p.strip() for p in notes.split('\n') if p.strip()]:
        F.append(Paragraph(esc(para), st_body))

def footer(canvas, docu):
    canvas.saveState()
    canvas.setFont('Helvetica', 8)
    canvas.setFillColor(MUTE)
    canvas.drawString(1.05*inch, 0.58*inch,
                      'HIST 201 · Week 1 · Kepler · lecture narration')
    canvas.drawRightString(LETTER[0]-1.05*inch, 0.58*inch, str(docu.page))
    canvas.restoreState()

doc.build(F, onFirstPage=footer, onLaterPages=footer)

words = sum(len(n.split()) for _, _, n in slides if n)
print('WROTE', OUT)
print('slides with narration:', sum(1 for _, _, n in slides if n))
print('total words:', words)
print('approx spoken minutes at 140 wpm: %.1f' % (words/140.0))
p1 = sum(len(n.split()) for i, _, n in slides if n and i < 7)
p2 = words - p1
print('  part one: %d words (%.1f min)' % (p1, p1/140.0))
print('  part two: %d words (%.1f min)' % (p2, p2/140.0))
